#!/usr/bin/env python3
"""Build the project defence slide deck (docs/AttendIQ_Defence.pptx).

Dark premium theme matching the app. Every slide carries speaker notes
so it doubles as a defence script. Re-run after doc updates:
    .venv/Scripts/python.exe scripts/build_defence_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x0F, 0x0F, 0x1A)
PANEL = RGBColor(0x1A, 0x1A, 0x28)
ACCENT = RGBColor(0xE9, 0x45, 0x60)   # warm red, app accent
COBALT = RGBColor(0x2E, 0x5A, 0xAC)
WHITE = RGBColor(0xF2, 0xF2, 0xF2)
GREY = RGBColor(0xA8, 0xA8, 0xB8)
GREEN = RGBColor(0x16, 0xC7, 0x9A)
AMBER = RGBColor(0xF5, 0xA6, 0x23)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def add_title(slide, text, size=40, color=WHITE, left=0.7, top=0.4, width=11.9):
    box = textbox(slide, left, top, width, 1.2)
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Calibri"
    return box


def add_accent_bar(slide, top=1.55, width=1.6):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(top), Inches(width), Pt(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_body(slide, lines, left=0.7, top=2.0, width=11.9, size=18, gap=6, bullets=True):
    box = textbox(slide, left, top, width, 4.5)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = 0
        if bullets:
            p.text = f"•  {line}"
        else:
            p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
    return box


def add_note(slide, text):
    slide.notes_slide.placeholders[1].text = text


def add_panel_row(slide, panels, top=2.1, height=3.4):
    n = len(panels)
    w = (11.9 - 0.5 * (n - 1)) / n
    for i, (title, body, color) in enumerate(panels):
        left = 0.7 + i * (w + 0.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL
        shape.line.fill.background()
        tbox = textbox(slide, left + 0.25, top + 0.2, w - 0.5, height - 0.4)
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        p2.text = body
        p2.font.size = Pt(15)
        p2.font.color.rgb = GREY
        p2.font.name = "Calibri"


def add_table_slide(slide, headers, rows, top=2.1):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.7), Inches(top), Inches(11.9), Inches(0.5))
    table = table_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = COBALT
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.color.rgb = WHITE
                p.font.name = "Calibri"
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PANEL


def footer(slide, text="AttendIQ  •  Final Year Defence  •  UNILAG Computer Engineering"):
    box = textbox(slide, 0.7, 7.0, 11.9, 0.4)
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = GREY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.RIGHT


# ---- 1. Title ----
s = prs.slides.add_slide(BLANK); bg(s)
add_title(s, "AttendIQ", size=72)
add_title(s, "Bias-Aware Facial Recognition Attendance System", size=30, color=ACCENT, top=1.7)
add_accent_bar(s, top=2.6, width=2.2)
box = textbox(s, 0.7, 3.0, 11.9, 2.0)
for line in ["Final Year Project Defence", "Department of Computer Engineering, University of Lagos", "2025/2026 Session"]:
    p = box.text_frame.add_paragraph() if len(box.text_frame.paragraphs[0].runs) else box.text_frame.paragraphs[0]
    p.text = line
    p.font.size = Pt(22)
    p.font.color.rgb = GREY
    p.font.name = "Calibri"
    p.space_after = Pt(4)
footer(s)
add_note(s, "Open with your name, matric number, supervisor. One line: manual attendance wastes class time and is easy to fake; commercial face systems are biased — AttendIQ tackles both and measures the bias.")

# ---- 2. Problem ----
s = prs.slides.add_slide(BLANK); bg(s)
add_title(s, "The problem", size=40); add_accent_bar(s)
add_body(s, [
    "Manual roll call burns 5–10 minutes of every lecture and proxy sign-ins are trivial.",
    "Fingerprint scanners need contact and queues; they do not scale with class size.",
    "Commercial face recognition is biased: Gender Shades found up to 34.7% higher error for dark-skinned users — a deployment without bias measurement is irresponsible.",
    "Supervisors need one system that is contactless, automatic, AND honest about where it underperforms.",
])
footer(s)
add_note(s, "Quantify: 5–10 min × 3 lectures/week × 15 weeks ≈ 4–7 hours lost per course. Then pivot: accuracy alone is not enough, fairness must be measured — that is the second half of the title.")

# ---- 3. Objectives ----
s = prs.slides.add_slide(BLANK); bg(s)
add_title(s, "Objectives", size=40); add_accent_bar(s)
add_body(s, [
    "Automate attendance with real-time face recognition (dlib 128-D primary, LBPH fallback).",
    "Support the full school workflow: enrollment, course registration, live sessions, manual correction, exports.",
    "Harden access with JWT + role-based control so attendance records cannot be spoofed with curl.",
    "Evaluate bias across Fitzpatrick skin types I–VI and gender using the Gender Shades method.",
    "Deploy on a headless VPS and a desktop shell from one codebase.",
])
footer(s)
add_note(s, "Five objectives map 1:1 to demo sections. If asked to cut scope, the bias module is the differentiator — never offer to drop it.")

# ---- 4. Architecture ----
s = prs.slides.add_slide(BLANK); bg(s)
add_title(s, "Architecture", size=40); add_accent_bar(s)
add_panel_row(s, [
    ("Browser client", "React 19 + TS + Tailwind. getUserMedia → canvas → Base64 JPEG. USB camera picker, role tabs, onboarding tour.", ACCENT),
    ("FastAPI backend", "REST + JWT (15m/7d) + RBAC + rate limits + security headers. Stateless frame recognition; MJPEG stream for local cameras.", COBALT),
    ("Vision engine", "Haar/DNN detection (equalizeHist, multi-face) → dlib 128-D or LBPH → Euclidean match at tolerance 0.6.", GREEN),
    ("Data + bias", "SQLite (users, classes, blocks, attendance). Gender Shades evaluator over skin type × gender.", AMBER),
], top=2.1, height=3.6)
footer(s)
add_note(s, "Walk left to right following one frame: browser captures, POST /api/recognize/frame, detection + encoding, match, attendance row, JSON boxes back. Mention headless-VPS design: server needs no camera.")

# ---- 5. Demo flow ----
s = prs.slides.add_slide(BLANK); bg(s)
add_title(s, "Live demo flow", size=40); add_accent_bar(s)
add_body(s, [
    "1. Lecturer signs in on the Lecturer tab — note the role-mismatch guard.",
    "2. Claim a course, open the roster: registered students, remove/block one → amber BLOCKED box.",
    "3. Start a live session: names marked in real time, chime on each new mark, presence dots catch walk-aways.",
    "4. Student view: register for a course with faculty/department filters; enrollment photo rules enforced.",
    "5. Admin view: bias charts, then export the session to CSV.",
], size=19)
footer(s)
add_note(s, "Do the demo in this order. Pre-enroll two faces and one blocked student before the defence. If the camera misbehaves, fall back to the browser-camera path — it uses the same backend.")

# ---- 6. Security ----
s = prs.slides.add_slide(BLANK); bg(s)
add_title(s, "Security: JWT + RBAC", size=40); add_accent_bar(s)
add_body(s, [
    "Old system: frontend-only guards — any student could POST /api/attendance/manual with curl and mark themselves present.",
    "Now: HS256 access (15 min) + rotating refresh (7d, hashed jti store); every route behind get_current_user / require_roles → 401 without token, 403 on role mismatch.",
    "Media fix: <img> tags cannot send headers, so the MJPEG feed accepts short-lived ?token= plus lecturer/admin check.",
    "Passwords PBKDF2-SHA256 (260k); security answers + emergency PINs hashed and invisible even to admins; parameterized SQL throughout.",
])
footer(s)
add_note(s, "If challenged, offer to demo the curl attack live: same request without token → 401, with a student token → 403. That usually ends the security questions.")

# ---- 7. Bias evaluation ----
s = prs.slides.add_slide(BLANK); bg(s)
add_title(s, "Bias evaluation", size=40); add_accent_bar(s)
add_body(s, [
    "Method: Gender Shades — accuracy sliced by Fitzpatrick types I–VI, gender, and their intersections.",
    "Metrics: detection rate, recognition accuracy, false-negative rate, disparity gap (best minus worst group).",
    "Pipeline: admin bootstraps data/evaluation_dataset, fills annotations.csv, runs evaluation; charts + metrics.json persist.",
    "Known pattern (NIST FRVT): lighter skin ≈99.7% vs darkest ≈94.6% — the module exists to surface our own gap, not hide it.",
])
footer(s)
add_note(s, "Supervisors love this slide. Stress that most student projects report accuracy only; yours reports accuracy per group plus the gap. Have your metrics.json numbers memorised.")

# ---- 8. Results ----
s = prs.slides.add_slide(BLANK); bg(s)
add_title(s, "Results & performance", size=40); add_accent_bar(s)
add_table_slide(s,
    ["Metric", "Value", "Condition"],
    [
        ["Detection rate", "> 95%", "Good front lighting"],
        ["Recognition accuracy", "> 90%", "5 enrolled photos, bare face"],
        ["False positives", "< 2%", "Tolerance 0.6 (dlib)"],
        ["End-to-end latency", "~30 ms / frame", "640×480, 0.25 scale"],
        ["Session dedup", "100%", "Set + UNIQUE constraint"],
    ])
footer(s)
add_note(s, "Replace these with YOUR measured numbers before defence day — run the evaluator on your dataset and edit this table. Never present placeholder benchmarks as your own.")

# ---- 9. Limitations ----
s = prs.slides.add_slide(BLANK); bg(s)
add_title(s, "Limitations (honest)", size=40); add_accent_bar(s)
add_body(s, [
    "Frontal faces: profiles and heavy occlusion still defeat Haar; extreme crowds (>6 faces at 640×480) drop distant faces.",
    "Software liveness only: blocks photos and screen replays, not 3D masks or deepfake video — needs IR/depth hardware.",
    "SQLite single-writer: fine for one department server, not for multi-server scale.",
    "No MFA yet; recovery is security-question/PIN based.",
    "PIR motion sensor and second camera are wired stubs, not live devices.",
], size=19)
footer(s)
add_note(s, "Own the limitations before they ask. Each one maps to a future-work item on the next slide — that turns an attack into a roadmap discussion.")

# ---- 10. Future work ----
s = prs.slides.add_slide(BLANK); bg(s)
add_title(s, "Future improvements", size=40); add_accent_bar(s)
add_body(s, [
    "Wire the PIR webhook + second camera into the streamer; auto-flag students unseen for a full session.",
    "Upgrade detection to MTCNN/RetinaFace and add IR/depth anti-spoofing.",
    "TOTP MFA, audit log, httponly Secure refresh cookies in production.",
    "Migrate to PostgreSQL (Supabase path: RLS + GoTrue, Storage for face BLOBs).",
    "CI/CD: GitHub Actions for tests, lint, and container builds; attendance trend analytics.",
], size=19)
footer(s)
add_note(s, "Pick the top two you could actually do in a month if asked. The hardware rollout is the most impressive-sounding and demos well with a cheap PIR module.")

# ---- 11. Conclusion ----
s = prs.slides.add_slide(BLANK); bg(s)
add_title(s, "Conclusion", size=44); add_accent_bar(s, width=2.4)
add_body(s, [
    "A working, deployed attendance system: ~3,950 lines Python + ~3,400 TypeScript, one codebase for VPS and desktop.",
    "Contributions: unified architecture, Gender Shades bias framework, JWT+RBAC hardening, multi-face + presence tracking.",
    "It does not just recognise faces — it reports where it fails, and refuses to log who should not be there.",
], size=20)
footer(s)
add_note(s, "Close in 30 seconds: problem, what you built, the one number that matters (your disparity gap), thank the supervisor and invite questions.")

# ---- 12. Q&A ----
s = prs.slides.add_slide(BLANK); bg(s)
add_title(s, "Questions?", size=72)
add_title(s, "Demo is live — ask me to break it.", size=26, color=ACCENT, top=2.2)
box = textbox(s, 0.7, 3.4, 11.9, 2.0)
for line in ["Likely questions: tolerance math  •  why dlib+LBPH  •  curl attack demo  •  blocked boxes  •  presence vs runaway  •  Supabase decision"]:
    p = box.text_frame.add_paragraph() if len(box.text_frame.paragraphs[0].runs) else box.text_frame.paragraphs[0]
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = GREY
    p.font.name = "Calibri"
footer(s)
add_note(s, "Have these ready: tolerance 0.6 Euclidean in 128-D; dlib accuracy vs LBPH zero-dependency fallback; curl 401/403 demo; BLOCKED boxes on both video paths; last_seen presence dots; Supabase deferred to keep single-VPS + desktop.")

prs.save("docs/AttendIQ_Defence.pptx")
print("Saved docs/AttendIQ_Defence.pptx")
